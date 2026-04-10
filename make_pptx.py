"""
DB_TO_LLM — 전체 아키텍처 & 워크플로우 발표용 PPT 생성 스크립트
슬라이드 구성 (13장):
 1  표지
 2  전체 시스템 개요
 3  프로젝트 폴더 구조
 4  공통 실행 구조
 5  Planner 워크플로우
 6  natural_llm 워크플로우
 7  prompt_llm 워크플로우
 8  rag_prompt_llm 워크플로우
 9  API_SERVER 워크플로우
10  End-to-End 시퀀스
11  실행 명령 / 운영 방법
12  현재 구현 범위 / 한계 / TODO
13  결론

실행: python make_pptx.py
"""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ──────────────────────────────────────────────
# 공통 색상 팔레트
# ──────────────────────────────────────────────
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK   = RGBColor(0x1A, 0x1A, 0x2E)
C_BLUE1   = RGBColor(0x00, 0x5F, 0xB5)   # 타이틀 계열
C_BLUE2   = RGBColor(0x00, 0x8D, 0xCE)    # 포인트 박스
C_ACCENT  = RGBColor(0x00, 0xB4, 0xD8)   # 화살표·강조
C_INGEST  = RGBColor(0x1B, 0x7F, 0x4D)   # Ingest 계열 (그린)
C_STREAM  = RGBColor(0x00, 0x5F, 0xB5)   # Stream 계열 (블루)
C_LITE    = RGBColor(0xE8, 0xF4, 0xFF)   # 밝은 배경용
C_GRAY    = RGBColor(0x55, 0x69, 0x7A)   # 보조 텍스트
C_YELLOW  = RGBColor(0xFF, 0xC3, 0x00)   # 경고·강조

# 슬라이드 크기 (와이드 16:9)
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

FONT_KO = "나눔고딕"  # 없으면 맑은 고딕으로 대체
FONT_FALLBACK = "맑은 고딕"


# ──────────────────────────────────────────────
# 헬퍼
# ──────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill: RGBColor | None = None,
             line_color: RGBColor | None = None, line_width: int = 1,
             radius: bool = False):
    """채워진 직사각형(rounded 옵션) 도형 추가."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1 if not radius else 5,  # 1=RECTANGLE, 5=ROUNDED_RECTANGLE
        x, y, w, h)
    shape.line.width = Pt(line_width)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text: str, x, y, w, h,
                font_size=12, bold=False, color: RGBColor = C_BLACK,
                align=PP_ALIGN.LEFT, font=None):
    """텍스트 박스 추가."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font or FONT_FALLBACK
    return txb


def add_arrow(slide, x1, y1, x2, y2, color=C_ACCENT, width=2):
    """라인 + 화살표 연결선."""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)   # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = Pt(width)


def set_slide_bg(slide, color: RGBColor = C_WHITE):
    """슬라이드 배경색 설정."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, text: str):
    """발표자 노트 삽입."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_title_bar(slide, title: str, subtitle: str | None = None,
                  bar_color=C_BLUE1):
    """상단 제목 바."""
    bar_h = Cm(2.6)
    add_rect(slide, Cm(0), Cm(0), SLIDE_W, bar_h, fill=bar_color)
    add_textbox(slide, title, Cm(0.7), Cm(0.2), SLIDE_W - Cm(1.4), Cm(1.3),
                font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, Cm(0.7), Cm(1.5), SLIDE_W - Cm(1.4), Cm(0.9),
                    font_size=12, bold=False, color=RGBColor(0xCC, 0xE5, 0xFF),
                    align=PP_ALIGN.LEFT)


def flow_box(slide, label: str, x, y, w=Cm(4.0), h=Cm(1.3),
             fill=C_BLUE2, text_color=C_WHITE, font_size=9, radius=True):
    """플로우 박스 + 텍스트."""
    add_rect(slide, x, y, w, h, fill=fill, radius=radius)
    add_textbox(slide, label, x, y + Cm(0.15), w, h,
                font_size=font_size, bold=True, color=text_color,
                align=PP_ALIGN.CENTER)


def arrow_h(slide, x_start, x_end, y, color=C_ACCENT):
    """수평 화살표."""
    add_arrow(slide, x_start, y, x_end, y, color=color)


def arrow_v(slide, x, y_start, y_end, color=C_ACCENT):
    """수직 화살표."""
    add_arrow(slide, x, y_start, x, y_end, color=color)


def add_key_message(slide, msg: str, y=None):
    """하단 핵심 메시지 바."""
    if y is None:
        y = SLIDE_H - Cm(1.6)
    add_rect(slide, Cm(0), y, SLIDE_W, Cm(1.5),
             fill=C_BLUE1)
    add_textbox(slide, f"💡 핵심 메시지: {msg}",
                Cm(0.7), y + Cm(0.15), SLIDE_W - Cm(1.4), Cm(1.2),
                font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

# ──────────────────────────────────────────────
# 슬라이드 생성 함수들
# ──────────────────────────────────────────────

# ───────── 슬라이드 01: 표지 ─────────
def slide_01_title(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)

    add_rect(slide, Cm(0), Cm(0), SLIDE_W, Cm(1.0), fill=C_BLUE1)
    add_rect(slide, Cm(1.5), Cm(1.8), SLIDE_W - Cm(3), Cm(8.5), fill=C_BLUE1, radius=True)

    add_textbox(slide, "DB_TO_LLM",
                Cm(2.0), Cm(2.5), SLIDE_W - Cm(4), Cm(2.8),
                font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "Planner / Stream / RAG / API Server\n아키텍처 및 워크플로우",
                Cm(2.0), Cm(5.6), SLIDE_W - Cm(4), Cm(2.2),
                font_size=18, bold=False, color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_rect(slide, Cm(7), Cm(7.9), Cm(20), Cm(0.08), fill=C_ACCENT)

    add_textbox(slide, "자연어 질문 → Planner 분기 → SQL/RAG/LLM → 최종 답변",
                Cm(2.0), Cm(8.2), SLIDE_W - Cm(4), Cm(0.9),
                font_size=11, color=RGBColor(0xCC, 0xE5, 0xFF), align=PP_ALIGN.CENTER)

    add_textbox(slide, "발표용  |  소스 기반 분석  |  2026",
                Cm(2.0), Cm(9.3), SLIDE_W - Cm(4), Cm(0.8),
                font_size=10, color=RGBColor(0xCC, 0xE5, 0xFF), align=PP_ALIGN.CENTER)

    add_rect(slide, Cm(0), Cm(0), Cm(0.4), SLIDE_H, fill=C_ACCENT)

    add_rect(slide, Cm(0), SLIDE_H - Cm(1.0), SLIDE_W, Cm(1.0), fill=C_BLUE1)
    add_textbox(slide, "Internal Use Only",
                Cm(1), SLIDE_H - Cm(0.9), SLIDE_W - Cm(2), Cm(0.8),
                font_size=9, color=C_WHITE, align=PP_ALIGN.RIGHT)

    add_notes(slide, "DB_TO_LLM 프로젝트의 전체 아키텍처와 5개 워크플로우(Planner, natural_llm, prompt_llm, rag_prompt_llm, API Server)를 한눈에 이해하기 위한 발표 자료입니다.")


# ───────── 슬라이드 02: 전체 시스템 개요 ─────────
def slide_02_system_overview(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "전체 시스템 개요",
                  "Root_Ingest / Root_Stream / Planner 관계 및 외부 의존성")

    # 세 영역: Planner(상단 우), Root_Ingest(좌), Root_Stream(우)
    # Root_Ingest 박스
    add_rect(slide, Cm(0.5), Cm(3.2), Cm(9.5), Cm(10.5),
             fill=RGBColor(0xE8, 0xF7, 0xEE), line_color=C_INGEST, radius=True)
    add_textbox(slide, "📥 Root_Ingest\n(오프라인 파이프라인)",
                Cm(0.8), Cm(3.4), Cm(9.0), Cm(1.0),
                font_size=11, bold=True, color=C_INGEST)
    ingest_steps = [
        "📂  문서 수집 (doc/)",
        "🔍  Parser / Loader",
        "✂️  Chunking (chunk_service)",
        "🧮  Embedding 생성",
        "🗄️  ChromaDB 적재",
    ]
    for i, s in enumerate(ingest_steps):
        flow_box(slide, s, Cm(0.9), Cm(4.6) + i * Cm(1.7),
                 w=Cm(9.0), h=Cm(1.3), fill=C_INGEST, font_size=8)
        if i < len(ingest_steps) - 1:
            arrow_v(slide, Cm(5.4), Cm(4.6) + i * Cm(1.7) + Cm(1.3),
                    Cm(4.6) + (i + 1) * Cm(1.7), color=C_INGEST)

    # Root_Stream 박스
    add_rect(slide, Cm(11.0), Cm(3.2), Cm(12.5), Cm(10.5),
             fill=C_LITE, line_color=C_BLUE1, radius=True)
    add_textbox(slide, "📤 Root_Stream\n(온라인 서빙)",
                Cm(11.3), Cm(3.4), Cm(12.0), Cm(1.0),
                font_size=11, bold=True, color=C_BLUE1)
    stream_steps = [
        "❓  자연어 질문 입력",
        "🔀  mode 분기 (Orchestrator)",
        "🔍  Vector 검색 (RAG 모드)",
        "🧩  Prompt 조합 (PromptManager)",
        "🤖  LLM 호출 (Ollama / OpenAI)",
    ]
    for i, s in enumerate(stream_steps):
        flow_box(slide, s, Cm(11.3), Cm(4.6) + i * Cm(1.7),
                 w=Cm(11.8), h=Cm(1.3), fill=C_BLUE1, font_size=8)
        if i < len(stream_steps) - 1:
            arrow_v(slide, Cm(17.2), Cm(4.6) + i * Cm(1.7) + Cm(1.3),
                    Cm(4.6) + (i + 1) * Cm(1.7), color=C_BLUE1)

    # Planner 박스 (상단 중앙)
    add_rect(slide, Cm(10.2), Cm(3.2), Cm(0.6), Cm(10.5),
             fill=RGBColor(0xDD, 0xDD, 0xDD))

    # ChromaDB 공유 화살표
    add_textbox(slide, "ChromaDB\n(공유)",
                Cm(9.7), Cm(7.2), Cm(1.8), Cm(1.3),
                font_size=8, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
    add_arrow(slide, Cm(10.0), Cm(8.0), Cm(10.0), Cm(6.2), color=C_YELLOW, width=2)
    add_arrow(slide, Cm(10.0), Cm(6.2), Cm(11.0), Cm(6.2), color=C_YELLOW, width=2)

    # Planner 버블 (우측 상단)
    add_rect(slide, Cm(24.5), Cm(3.2), Cm(9.0), Cm(4.5),
             fill=RGBColor(0xFF, 0xF3, 0xE0), line_color=RGBColor(0xFF, 0x8F, 0x00), radius=True)
    add_textbox(slide, "🧠 Planner",
                Cm(24.8), Cm(3.4), Cm(8.5), Cm(0.8),
                font_size=11, bold=True, color=RGBColor(0xFF, 0x8F, 0x00))
    planner_items = [
        "질문 분석 → query_type 결정",
        "DB_ONLY / RAG_ONLY / GENERAL",
        "DB_THEN_RAG / DB_THEN_GENERAL",
        "RAG_THEN_GENERAL",
    ]
    for i, s in enumerate(planner_items):
        add_textbox(slide, f"• {s}", Cm(25.0), Cm(4.4) + i * Cm(0.7), Cm(8.0), Cm(0.65),
                    font_size=8, color=RGBColor(0x55, 0x40, 0x00))

    # 외부 의존성 박스
    add_rect(slide, Cm(24.5), Cm(8.2), Cm(9.0), Cm(5.5),
             fill=RGBColor(0xF3, 0xEE, 0xFF), line_color=RGBColor(0x6A, 0x00, 0x8E), radius=True)
    add_textbox(slide, "🔌 외부 의존성",
                Cm(24.8), Cm(8.4), Cm(8.5), Cm(0.7),
                font_size=10, bold=True, color=RGBColor(0x6A, 0x00, 0x8E))
    ext_deps = [
        ("LLM", "Ollama (qwen2.5:14b) / OpenAI"),
        ("VectorDB", "ChromaDB PersistentClient"),
        ("RDBMS", "MSSQL (pymssql)"),
        ("Server", "FastAPI (uvicorn)"),
        ("임베딩", "sentence-transformers (multilingual)"),
    ]
    for i, (k, v) in enumerate(ext_deps):
        add_textbox(slide, f"• {k}: {v}",
                    Cm(25.0), Cm(9.2) + i * Cm(0.8), Cm(8.0), Cm(0.7),
                    font_size=8, color=RGBColor(0x40, 0x00, 0x60))

    add_key_message(slide, "Root_Ingest가 지식을 준비하고, Root_Stream이 실시간으로 응답한다. Planner는 질문 유형을 분류해 흐름을 결정한다.")
    add_notes(slide, "Root_Ingest(오프라인)와 Root_Stream(온라인)이 ChromaDB를 공유합니다. Planner는 질문을 분석해 query_type을 결정하고, Root_Stream이 해당 유형에 맞게 SQL/RAG/LLM을 실행합니다.")


# ───────── 슬라이드 03: 폴더 구조 ─────────
def slide_03_folder_structure(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "프로젝트 폴더 구조 요약",
                  "각 모듈의 역할과 책임 범위")

    folders = [
        {
            "name": "Planner/",
            "color": RGBColor(0xFF, 0x8F, 0x00),
            "bg": RGBColor(0xFF, 0xF3, 0xE0),
            "desc": "질문 분석 & 실행 계획 수립",
            "files": [
                "planner_service.py  — PlannerService 진입점",
                "models.py           — PlannerPlan / PlannerStep dataclass",
                "planner_prompt.py   — prompt_templates.yaml 키 바인딩",
                "plan_validator.py   — JSON 구조 최소 검증",
                "debug_planner.py    — CLI 디버그 실행",
            ],
        },
        {
            "name": "Root_Ingest/",
            "color": C_INGEST,
            "bg": RGBColor(0xE8, 0xF7, 0xEE),
            "desc": "문서 수집 → 파싱 → 청킹 → 임베딩 → ChromaDB 적재",
            "files": [
                "ingest/ingest_pipeline.py  — 전체 파이프라인 오케스트레이터",
                "ingest/document_loader.py  — 파일 수집·인덱싱",
                "ingest/parser_service.py   — SQL DDL / 문서 파싱",
                "ingest/chunk_service.py    — 청킹 (size=800, overlap=100)",
                "ingest/embedding_service.py — SentenceTransformer 임베딩",
                "ingest/vector_store_service.py — ChromaDB upsert",
            ],
        },
        {
            "name": "Root_Stream/",
            "color": C_BLUE1,
            "bg": C_LITE,
            "desc": "자연어 질문 → SQL 생성 → 응답 반환 (3가지 mode)",
            "files": [
                "orchestrator/stream_orchestrator.py — mode 분기 총괄",
                "services/stream/mode_natural_llm.py — 직접 LLM 호출",
                "services/stream/mode_prompt_llm.py  — 템플릿 기반 SQL 생성",
                "services/stream/mode_rag_prompt_llm.py — RAG + LLM",
                "services/sql/sql_guard.py           — SELECT-only 안전 검증",
                "services/sql/sql_executor_service.py — MSSQL 실행",
                "services/e2e/planner_flow_runner.py — E2E 흐름 오케스트레이션",
                "server/api_app.py + routes.py       — FastAPI 서버",
                "prompts/prompt_templates.yaml       — 프롬프트 템플릿 저장소",
            ],
        },
        {
            "name": "docs/ & tests/",
            "color": C_GRAY,
            "bg": RGBColor(0xF5, 0xF5, 0xF5),
            "desc": "SQL 스키마 문서 및 단위·통합 테스트",
            "files": [
                "doc/ALL_LOG_TABLES_cleaned.sql — DB 스키마 DDL",
                "tests/root_stream/             — config_loader, SQL guard 테스트",
            ],
        },
    ]

    col_w = (SLIDE_W - Cm(1.2)) / 2 - Cm(0.3)
    positions = [(Cm(0.5), Cm(3.2)), (Cm(0.5 + col_w + 0.6), Cm(3.2)),
                 (Cm(0.5), Cm(10.0)), (Cm(0.5 + col_w + 0.6), Cm(10.0))]

    for (gx, gy), folder in zip(positions, folders):
        box_h = Cm(6.5)
        add_rect(slide, gx, gy, col_w, box_h,
                 fill=folder["bg"], line_color=folder["color"], radius=True)
        # 헤더
        add_rect(slide, gx, gy, col_w, Cm(1.1), fill=folder["color"], radius=True)
        add_textbox(slide, folder["name"],
                    gx + Cm(0.3), gy + Cm(0.1), col_w - Cm(0.4), Cm(0.6),
                    font_size=11, bold=True, color=C_WHITE)
        add_textbox(slide, folder["desc"],
                    gx + Cm(0.3), gy + Cm(0.7), col_w - Cm(0.4), Cm(0.5),
                    font_size=8, color=folder["color"])
        for i, f in enumerate(folder["files"]):
            add_textbox(slide, f"  {f}",
                        gx + Cm(0.2), gy + Cm(1.3) + i * Cm(0.78), col_w - Cm(0.3), Cm(0.72),
                        font_size=7.5, color=RGBColor(0x33, 0x33, 0x33))

    add_key_message(slide, "폴더 = 책임 단위. Planner/Root_Ingest/Root_Stream이 독립적으로 분리되어 있고, Root_Stream이 나머지를 통합 호출한다.")
    add_notes(slide, "각 폴더는 책임이 명확히 분리됩니다. Planner와 Root_Ingest는 Root_Stream이 임포트해 사용하는 의존성입니다.")


# ───────── 슬라이드 04: 공통 실행 구조 ─────────
def slide_04_common_structure(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "공통 실행 구조",
                  "모든 실행 경로가 공유하는 공통 컴포넌트 관계도")

    # 공통 컴포넌트 레이아웃 (가로 배치)
    components = [
        ("config.yaml\n+ .env",       "설정 주입\n(mode, provider, paths)",   RGBColor(0x5C, 0x6B, 0xC0)),
        ("prompt_templates\n.yaml",   "프롬프트 키·템플릿\n저장소",           RGBColor(0x5C, 0x6B, 0xC0)),
        ("PromptManager",             "get_prompt(key)\nrender_prompt(key,{})", C_BLUE1),
        ("LLM Client\n(Factory)",     "OllamaClient\nOpenAIClient",            RGBColor(0x8B, 0x00, 0x00)),
        ("Retrieval\n(RAG 전용)",     "SentenceTransformer\n+ ChromaRetriever", C_INGEST),
        ("SQL Guard",                 "SELECT-only 검증\n금지 키워드 차단",     RGBColor(0xE6, 0x37, 0x00)),
        ("DB Execution\n(선택적)",    "SqlExecutorService\n→ MssqlClient",      RGBColor(0x55, 0x55, 0x55)),
    ]

    comp_w = Cm(4.1)
    comp_h = Cm(2.8)
    total_w = len(components) * comp_w + (len(components) - 1) * Cm(0.5)
    start_x = (SLIDE_W - total_w) / 2
    box_y = Cm(3.5)

    centers = []
    for i, (title, desc, col) in enumerate(components):
        bx = start_x + i * (comp_w + Cm(0.5))
        add_rect(slide, bx, box_y, comp_w, comp_h,
                 fill=col, line_color=col, radius=True)
        add_textbox(slide, title, bx, box_y + Cm(0.2), comp_w, Cm(1.0),
                    font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, bx + Cm(0.1), box_y + Cm(1.3), comp_w - Cm(0.2), Cm(1.3),
                 fill=RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, desc, bx + Cm(0.15), box_y + Cm(1.35), comp_w - Cm(0.3), Cm(1.2),
                    font_size=7.5, color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.CENTER)
        cx = bx + comp_w / 2
        centers.append(cx)
        if i < len(components) - 1:
            arrow_h(slide, bx + comp_w, bx + comp_w + Cm(0.5), box_y + comp_h / 2, color=C_GRAY)

    # 관계 설명 박스들
    relations = [
        ("load_config(config.yaml)", C_BLUE1),
        ("→ PromptManager(prompt_file)", RGBColor(0x5C, 0x6B, 0xC0)),
        ("→ create_llm_client(config)", RGBColor(0x8B, 0x00, 0x00)),
        ("→ (RAG) embed + retrieve", C_INGEST),
        ("→ SqlGuard.validate()", RGBColor(0xE6, 0x37, 0x00)),
        ("→ SqlExecutorService.execute()", RGBColor(0x55, 0x55, 0x55)),
    ]
    add_textbox(slide, "공통 초기화 흐름:",
                Cm(0.7), Cm(7.0), Cm(5.0), Cm(0.6),
                font_size=10, bold=True, color=C_BLUE1)
    for i, (rel_txt, col) in enumerate(relations):
        add_textbox(slide, rel_txt,
                    Cm(0.9), Cm(7.7) + i * Cm(0.85), Cm(18.0), Cm(0.8),
                    font_size=9, color=col)

    # config.yaml 주요 키 설명
    add_rect(slide, Cm(20.0), Cm(7.0), Cm(13.3), Cm(9.0),
             fill=RGBColor(0xF8, 0xF8, 0xFF), line_color=C_BLUE1, radius=True)
    add_textbox(slide, "config.yaml 주요 설정",
                Cm(20.3), Cm(7.2), Cm(13.0), Cm(0.7),
                font_size=10, bold=True, color=C_BLUE1)
    cfg_items = [
        "mode: natural_llm | prompt_llm | rag_prompt_llm",
        "llm_provider: ollama | openai",
        "ollama.model: qwen2.5:14b",
        "retrieval.enabled: true",
        "retrieval.chroma_path: ../Root_Ingest/data/chroma",
        "retrieval.collection_name: doc_chunks",
        "retrieval.top_k: 3",
        "prompts.active_prompt: query_generation_prompt",
        "database.host / .database / .port",
        "sql.allow_only_select: true",
    ]
    for i, item in enumerate(cfg_items):
        add_textbox(slide, f"  {item}",
                    Cm(20.3), Cm(8.0) + i * Cm(0.76), Cm(12.8), Cm(0.72),
                    font_size=8, color=RGBColor(0x33, 0x33, 0x33))

    add_key_message(slide, "모든 모드는 config.yaml → PromptManager → LLMClient 초기화를 동일하게 거친다. 차이점은 Retrieval과 DB 실행 여부뿐이다.")
    add_notes(slide, "config.yaml이 모든 실행 경로의 중심입니다. mode 값으로 natural_llm/prompt_llm/rag_prompt_llm을 선택하고, Planner는 query_type으로 DB/RAG/LLM 경로를 추가로 결정합니다.")


# ───────── 슬라이드 05: Planner 워크플로우 ─────────
def slide_05_planner(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "Planner 워크플로우",
                  "질문 → JSON 실행 계획 생성 → query_type 결정")

    planner_color = RGBColor(0xFF, 0x8F, 0x00)
    planner_bg = RGBColor(0xFF, 0xF3, 0xE0)
    done_color = C_INGEST
    todo_color = C_GRAY

    # 메인 플로우 (수평, 4개 주요 단계)
    steps = [
        ("① 질문 입력\n(question)", "run_planner_step()\n← 04_end_to_end\n  planner_flow.ipynb", planner_color),
        ("② Planner 프롬프트\n조합", "planner_prompt.py\nbuild_planner_prompts()\nplanner_system_prompt\n+ planner_user_prompt", planner_color),
        ("③ LLM 호출\n(온도=0.0)", "PlannerService\n.plan_question()\ncreate_llm_client(config)\n→ generate()", planner_color),
        ("④ JSON 파싱\n& 형식 검증", "PlannerJsonParseError\nplan_validator.py\nvalidate_plan_payload()\nPlannerPlan.from_dict()", done_color),
    ]

    step_w, step_h = Cm(7.5), Cm(4.0)
    total_w = len(steps) * step_w + (len(steps) - 1) * Cm(0.8)
    sx = (SLIDE_W - total_w) / 2
    sy = Cm(3.3)

    for i, (title, desc, col) in enumerate(steps):
        bx = sx + i * (step_w + Cm(0.8))
        add_rect(slide, bx, sy, step_w, step_h, fill=col, radius=True)
        add_textbox(slide, title, bx, sy + Cm(0.2), step_w, Cm(1.0),
                    font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, bx + Cm(0.1), sy + Cm(1.2), step_w - Cm(0.2), Cm(2.5),
                 fill=C_WHITE)
        add_textbox(slide, desc, bx + Cm(0.2), sy + Cm(1.3), step_w - Cm(0.3), Cm(2.3),
                    font_size=7.5, color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.LEFT)
        if i < len(steps) - 1:
            arrow_h(slide, bx + step_w, bx + step_w + Cm(0.8), sy + step_h / 2, color=planner_color)

    # 출력: query_type 분기
    add_rect(slide, Cm(0.5), Cm(8.3), SLIDE_W - Cm(1.0), Cm(3.8),
             fill=planner_bg, line_color=planner_color, radius=True)
    add_textbox(slide, "Planner 출력: query_type 분기",
                Cm(0.8), Cm(8.5), Cm(10.0), Cm(0.7),
                font_size=10, bold=True, color=planner_color)

    qt_items = [
        ("DB_ONLY", "SQL 생성 + DB 실행만 수행", done_color),
        ("RAG_ONLY", "문서 RAG 검색 중심으로 답변", done_color),
        ("GENERAL", "DB/RAG 없이 LLM 직접 답변", done_color),
        ("DB_THEN_RAG", "DB 결과 → 문서 RAG 추가 수행", done_color),
        ("DB_THEN_GENERAL", "DB 결과 → 일반 설명 마무리", done_color),
        ("RAG_THEN_GENERAL", "문서 RAG 먼저 → 일반 설명 마무리", done_color),
    ]
    qt_w = (SLIDE_W - Cm(2.0)) / len(qt_items) - Cm(0.3)
    for i, (qt, desc, col) in enumerate(qt_items):
        gx = Cm(0.8) + i * (qt_w + Cm(0.3))
        flow_box(slide, qt, gx, Cm(9.3), w=qt_w, h=Cm(0.9), fill=col, font_size=8)
        add_textbox(slide, desc, gx, Cm(10.3), qt_w, Cm(0.8),
                    font_size=7, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 현재 구현 / 미구현 표시
    add_rect(slide, Cm(0.5), Cm(12.4), Cm(16.0), Cm(4.7),
             fill=RGBColor(0xE8, 0xF7, 0xEE), line_color=done_color, radius=True)
    add_textbox(slide, "✅ 현재 구현됨",
                Cm(0.8), Cm(12.6), Cm(15.5), Cm(0.6),
                font_size=9, bold=True, color=done_color)
    impl_items = [
        "PlannerService.plan_question() — LLM → JSON 파싱 → 검증",
        "PlannerPlan / PlannerStep dataclass (models.py)",
        "plan_validator.py — query_type, steps 최소 구조 검증",
        "planner_prompt.py — prompt_templates.yaml 키 바인딩",
        "planner_flow_runner.py — run_planner_step() E2E 연계",
    ]
    for i, item in enumerate(impl_items):
        add_textbox(slide, f"  • {item}",
                    Cm(0.9), Cm(13.3) + i * Cm(0.7), Cm(15.5), Cm(0.65),
                    font_size=8, color=RGBColor(0x1B, 0x5E, 0x20))

    add_rect(slide, Cm(17.0), Cm(12.4), Cm(16.3), Cm(4.7),
             fill=RGBColor(0xFF, 0xF3, 0xE0), line_color=todo_color, radius=True)
    add_textbox(slide, "⚠️ 현재 제외 / TODO",
                Cm(17.3), Cm(12.6), Cm(16.0), Cm(0.6),
                font_size=9, bold=True, color=todo_color)
    todo_items = [
        "Planner 독립 API 엔드포인트 — 미구현",
        "Planner 결과의 자동 캐싱 — 미구현",
        "query_type별 자동 라우팅 직접 연결 — 별도 E2E 함수로만 구현",
        "streaming 응답 지원 — 미구현",
    ]
    for i, item in enumerate(todo_items):
        add_textbox(slide, f"  • {item}",
                    Cm(17.4), Cm(13.3) + i * Cm(0.7), Cm(16.0), Cm(0.65),
                    font_size=8, color=RGBColor(0x99, 0x60, 0x00))

    add_key_message(slide, "Planner는 질문을 JSON 실행 계획(query_type + steps 배열)으로 변환한다. DB/RAG/LLM 중 어떤 경로를 쓸지 결정하는 라우터 역할이다.")
    add_notes(slide, "Planner는 Planner/planner_service.py의 PlannerService.plan_question()이 핵심입니다. LLM을 호출해 JSON을 받고, plan_validator로 검증 후 PlannerPlan으로 변환합니다. DB/RAG/LLM 실행은 포함하지 않습니다.")


# ───────── 슬라이드 06: natural_llm 워크플로우 ─────────
def slide_06_natural_llm(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "natural_llm 워크플로우",
                  "최소 처리 — 질문을 LLM에 직접 전달해 SQL 생성")

    nat_color = RGBColor(0x3A, 0x86, 0xFF)
    bg_color = RGBColor(0xEE, 0xF4, 0xFF)

    # 플로우 다이어그램 (수직 + 설명 오른쪽)
    steps_data = [
        ("① 사용자 질문 입력",
         "StreamRequest(question=...) 생성",
         "stream_orchestrator.py → run(question)"),
        ("② mode 확인",
         "config.mode = 'natural_llm'",
         "StreamOrchestrator._resolve_mode() → 'natural_llm'"),
        ("③ default_system_prompt 로딩",
         "prompt_manager.get_prompt('default_system_prompt')",
         "MSSQL 전용 Text-to-SQL 전문가 지시문"),
        ("④ user_prompt = question.strip()",
         "별도 템플릿 없음 — 질문 그대로 전달",
         "schema_context / business_rules 없음"),
        ("⑤ LLM 호출",
         "llm_client.generate(system_prompt, user_prompt)",
         "Ollama(qwen2.5:14b) 또는 OpenAI"),
        ("⑥ StreamResult 반환",
         "mode='natural_llm', query=생성된SQL",
         "CLI → JSON 출력 / FastAPI → HTTP 200"),
    ]

    box_w, box_h = Cm(4.5), Cm(1.1)
    desc_x = Cm(5.5)
    desc_w = Cm(13.0)
    src_x = Cm(19.0)
    src_w = Cm(14.3)
    sy = Cm(3.3)
    gap = Cm(1.9)

    for i, (label, desc, src) in enumerate(steps_data):
        yy = sy + i * gap
        flow_box(slide, label, Cm(0.5), yy, w=box_w, h=box_h, fill=nat_color, font_size=8)
        add_textbox(slide, desc, desc_x, yy + Cm(0.1), desc_w, box_h,
                    font_size=9, bold=True, color=C_BLUE1)
        add_textbox(slide, src, src_x, yy + Cm(0.1), src_w, box_h,
                    font_size=8, color=C_GRAY)
        if i < len(steps_data) - 1:
            arrow_v(slide, Cm(0.5) + box_w / 2,
                    yy + box_h, yy + box_h + Cm(0.8), color=nat_color)

    # 특징과 한계 박스
    add_rect(slide, Cm(0.5), Cm(15.3), Cm(15.5), Cm(3.0),
             fill=bg_color, line_color=nat_color, radius=True)
    add_textbox(slide, "✅  natural_llm 특징",
                Cm(0.8), Cm(15.5), Cm(15.0), Cm(0.6),
                font_size=9, bold=True, color=nat_color)
    nat_pros = [
        "• 구현 단순 — 프롬프트 키 'default_system_prompt' 하나만 사용",
        "• 추가 데이터 없이 즉시 실행 가능 (ChromaDB 불필요)",
        "• 테스트/프로토타이핑에 적합",
    ]
    for i, s in enumerate(nat_pros):
        add_textbox(slide, s, Cm(0.9), Cm(16.2) + i * Cm(0.65), Cm(15.0), Cm(0.6),
                    font_size=8, color=RGBColor(0x1B, 0x5E, 0x20))

    add_rect(slide, Cm(17.0), Cm(15.3), Cm(16.3), Cm(3.0),
             fill=RGBColor(0xFF, 0xF3, 0xE0), line_color=C_GRAY, radius=True)
    add_textbox(slide, "⚠️  한계",
                Cm(17.3), Cm(15.5), Cm(16.0), Cm(0.6),
                font_size=9, bold=True, color=C_GRAY)
    nat_cons = [
        "• 스키마 컨텍스트 없음 → 없는 테이블명 생성 가능성",
        "• RAG 검색 없음 → 도메인 지식 활용 불가",
        "• 실제 운영보다는 빠른 확인용으로 권장",
    ]
    for i, s in enumerate(nat_cons):
        add_textbox(slide, s, Cm(17.4), Cm(16.2) + i * Cm(0.65), Cm(16.0), Cm(0.6),
                    font_size=8, color=C_GRAY)

    add_key_message(slide, "natural_llm = 가장 단순한 경로. 스키마 없이 LLM에 질문을 그대로 전달한다. 빠른 테스트 용도에 적합하다.")
    add_notes(slide, "natural_llm은 Root_Stream/services/stream/mode_natural_llm.py에 구현됩니다. system_prompt는 MSSQL Text-to-SQL 지시문이고, user_prompt는 사용자 질문 그대로입니다. schema_context가 없으므로 정확도 제한이 있습니다.")


# ───────── 슬라이드 07: prompt_llm 워크플로우 ─────────
def slide_07_prompt_llm(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "prompt_llm 워크플로우",
                  "스키마·업무 규칙을 템플릿에 결합해 SQL 생성 품질 향상")

    prompt_color = RGBColor(0x6A, 0x00, 0x8E)
    prompt_bg = RGBColor(0xF3, 0xEE, 0xFF)

    steps_data = [
        ("① 질문 입력",
         "StreamRequest(question=...)",
         "stream_orchestrator.py → mode='prompt_llm'"),
        ("② active_prompt 키 읽기",
         "config.prompts.active_prompt\n= 'query_generation_prompt'",
         "mode_prompt_llm.py\n→ prompts_config.get('active_prompt')"),
        ("③ 템플릿 변수 구성",
         "_build_prompt_values() 호출",
         "question / schema_context\nbusiness_rules / additional_constraints"),
        ("④ 프롬프트 렌더링",
         "prompt_manager.render_prompt(\n  'query_generation_prompt', values)",
         "prompt_templates.yaml 내\n{question} {schema_context} 치환"),
        ("⑤ LLM 호출",
         "llm_client.generate(\n  system_prompt, user_prompt)",
         "Ollama(qwen2.5:14b) / OpenAI"),
        ("⑥ StreamResult 반환",
         "mode='prompt_llm'\nprompt_key='query_generation_prompt'",
         "CLI → JSON / FastAPI → HTTP 200"),
    ]

    box_w, box_h = Cm(4.5), Cm(1.1)
    desc_x = Cm(5.5)
    desc_w = Cm(13.0)
    src_x = Cm(19.0)
    src_w = Cm(14.3)
    sy = Cm(3.3)
    gap = Cm(1.85)

    for i, (label, desc, src) in enumerate(steps_data):
        yy = sy + i * gap
        flow_box(slide, label, Cm(0.5), yy, w=box_w, h=box_h, fill=prompt_color, font_size=8)
        add_textbox(slide, desc, desc_x, yy + Cm(0.05), desc_w, box_h + Cm(0.2),
                    font_size=8.5, bold=False, color=prompt_color)
        add_textbox(slide, src, src_x, yy + Cm(0.05), src_w, box_h + Cm(0.2),
                    font_size=8, color=C_GRAY)
        if i < len(steps_data) - 1:
            arrow_v(slide, Cm(0.5) + box_w / 2,
                    yy + box_h, yy + box_h + Cm(0.75), color=prompt_color)

    # 템플릿 변수 설명 박스
    add_rect(slide, Cm(0.5), Cm(15.0), Cm(15.5), Cm(3.5),
             fill=prompt_bg, line_color=prompt_color, radius=True)
    add_textbox(slide, "query_generation_prompt 변수",
                Cm(0.8), Cm(15.2), Cm(15.0), Cm(0.6),
                font_size=9, bold=True, color=prompt_color)
    tpl_vars = [
        ("{question}", "사용자 자연어 질문"),
        ("{schema_context}", "ERROR_LOG_DATA / EVENT_LOG_DATA / WARNING_LOG_DATA 테이블 설명"),
        ("{business_rules}", "분석 규칙 (EQPID 기준, ALARMID+TEXT 함께 반환 등)"),
        ("{additional_constraints}", "추가 제약 (기본값: 빈 문자열)"),
    ]
    for i, (var, desc) in enumerate(tpl_vars):
        add_textbox(slide, f"  {var}", Cm(0.9), Cm(15.9) + i * Cm(0.65), Cm(5.5), Cm(0.6),
                    font_size=8.5, bold=True, color=prompt_color)
        add_textbox(slide, desc, Cm(6.5), Cm(15.9) + i * Cm(0.65), Cm(9.0), Cm(0.6),
                    font_size=8, color=C_GRAY)

    add_rect(slide, Cm(17.0), Cm(15.0), Cm(16.3), Cm(3.5),
             fill=RGBColor(0xE8, 0xF7, 0xEE), line_color=C_INGEST, radius=True)
    add_textbox(slide, "✅  natural_llm 대비 개선점",
                Cm(17.3), Cm(15.2), Cm(16.0), Cm(0.6),
                font_size=9, bold=True, color=C_INGEST)
    pros = [
        "• schema_context 포함 → 실제 테이블/컬럼명 기반 SQL 생성",
        "• business_rules 포함 → 도메인 분석 규칙 반영",
        "• RAG 없이도 높은 정확도 가능 (스키마 명시적 제공 시)",
        "• ChromaDB 불필요 — 빠른 응답",
    ]
    for i, s in enumerate(pros):
        add_textbox(slide, s, Cm(17.4), Cm(15.9) + i * Cm(0.65), Cm(16.0), Cm(0.6),
                    font_size=8, color=RGBColor(0x1B, 0x5E, 0x20))

    add_key_message(slide, "prompt_llm = 스키마·업무 규칙을 프롬프트 템플릿에 주입해 SQL 품질을 높인다. RAG 없이 가장 현실적인 중간 단계이다.")
    add_notes(slide, "prompt_llm은 Root_Stream/services/stream/mode_prompt_llm.py에 구현됩니다. config.yaml의 prompts 섹션에서 schema_context와 business_rules를 읽어 템플릿에 주입합니다.")


# ───────── 슬라이드 08: rag_prompt_llm 워크플로우 ─────────
def slide_08_rag_prompt_llm(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "rag_prompt_llm 워크플로우",
                  "ChromaDB 검색 결과를 프롬프트에 결합 → 가장 높은 SQL 정확도")

    rag_color = RGBColor(0x00, 0x5F, 0xB5)
    ingest_link_color = C_INGEST

    steps_data = [
        ("① 질문 입력",
         "StreamRequest(question=...)",
         "mode='rag_prompt_llm'\nretrieval.enabled=true 필요"),
        ("② 질의 임베딩",
         "SentenceTransformerEmbeddingService\n.embed_query(question)\n→ list[float]",
         "embedding_model:\nparaphrase-multilingual\n-MiniLM-L12-v2"),
        ("③ ChromaDB 검색",
         "ChromaRetriever.retrieve(query_embedding)\ncollection='doc_chunks', top_k=3",
         "← Root_Ingest가 적재한\nChromaDB 컬렉션 사용"),
        ("④ context 블록 구성",
         "_build_context_block(contexts)\n검색된 청크 텍스트 결합",
         "RetrievedContext 리스트\n(chunk_id, text, score, metadata)"),
        ("⑤ 프롬프트 렌더링",
         "prompt_manager.render_prompt(\n  'rag_query_generation_prompt', values)",
         "{question} + {retrieved_context}\n+ schema_context + constraints"),
        ("⑥ LLM 호출",
         "llm_client.generate(system, user)\n→ SQL 반환",
         "Ollama(qwen2.5:14b) / OpenAI"),
        ("⑦ StreamResult 반환",
         "mode='rag_prompt_llm'\nretrieved_contexts=[...] 포함",
         "검색 근거 추적 가능\n(디버깅·설명 가능성↑)"),
    ]

    box_w, box_h = Cm(4.2), Cm(1.1)
    desc_x = Cm(5.0)
    desc_w = Cm(13.5)
    src_x = Cm(18.8)
    src_w = Cm(14.5)
    sy = Cm(3.3)
    gap = Cm(1.72)

    for i, (label, desc, src) in enumerate(steps_data):
        yy = sy + i * gap
        col = ingest_link_color if i == 2 else rag_color
        flow_box(slide, label, Cm(0.5), yy, w=box_w, h=box_h, fill=col, font_size=8)
        add_textbox(slide, desc, desc_x, yy + Cm(0.02), desc_w, box_h + Cm(0.2),
                    font_size=8.2, color=rag_color)
        add_textbox(slide, src, src_x, yy + Cm(0.02), src_w, box_h + Cm(0.2),
                    font_size=7.8, color=C_GRAY)
        if i < len(steps_data) - 1:
            arrow_v(slide, Cm(0.5) + box_w / 2,
                    yy + box_h, yy + box_h + Cm(0.62), color=col)

    # Root_Ingest 연결 강조 박스
    add_rect(slide, Cm(0.5), Cm(15.7), Cm(15.5), Cm(2.5),
             fill=RGBColor(0xE8, 0xF7, 0xEE), line_color=ingest_link_color, radius=True)
    add_textbox(slide, "Root_Ingest 연결 관계",
                Cm(0.8), Cm(15.9), Cm(15.0), Cm(0.6),
                font_size=9, bold=True, color=ingest_link_color)
    link_items = [
        "Root_Ingest/data/chroma  ──▶  retrieval.chroma_path 설정으로 rag_prompt_llm이 직접 읽음",
        "Root_Ingest가 생성한 collection_name='doc_chunks'를 collection_name 설정으로 지정",
        "임베딩 모델 일치 필수: Root_Ingest embedding.model_name == retrieval.embedding_model",
    ]
    for i, item in enumerate(link_items):
        add_textbox(slide, f"  • {item}",
                    Cm(0.9), Cm(16.6) + i * Cm(0.55), Cm(15.0), Cm(0.52),
                    font_size=7.8, color=RGBColor(0x1B, 0x5E, 0x20))

    # 성능 특징
    add_rect(slide, Cm(17.0), Cm(15.7), Cm(16.3), Cm(2.5),
             fill=RGBColor(0xE3, 0xF2, 0xFD), line_color=rag_color, radius=True)
    add_textbox(slide, "✅  rag_prompt_llm 특징",
                Cm(17.3), Cm(15.9), Cm(16.0), Cm(0.6),
                font_size=9, bold=True, color=rag_color)
    rag_pros = [
        "• ChromaDB top-k 청크로 테이블 DDL·설명 동적 제공",
        "• 동적 컨텍스트 → 스키마 업데이트 시 자동 반영",
        "• retrieved_contexts 포함 → 응답 근거 추적 가능",
    ]
    for i, s in enumerate(rag_pros):
        add_textbox(slide, s, Cm(17.4), Cm(16.6) + i * Cm(0.55), Cm(16.0), Cm(0.52),
                    font_size=7.8, color=rag_color)

    add_key_message(slide, "rag_prompt_llm = Root_Ingest가 만든 ChromaDB를 검색해 동적 컨텍스트를 주입한다. 현재 config.yaml의 기본 mode이다.")
    add_notes(slide, "rag_prompt_llm은 Root_Stream/services/stream/mode_rag_prompt_llm.py에 구현됩니다. retrieval.enabled=true 필요. Root_Ingest가 적재한 ChromaDB 컬렉션을 직접 읽습니다.")


# ───────── 슬라이드 09: API_SERVER 워크플로우 ─────────
def slide_09_api_server(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "API_SERVER 워크플로우",
                  "FastAPI 서버 — HTTP 요청 → mode 처리 → SQL 생성 응답")

    api_color = RGBColor(0x00, 0x96, 0x88)
    api_bg = RGBColor(0xE0, 0xF7, 0xFA)

    # 좌측: 서버 컴포넌트 구조
    add_rect(slide, Cm(0.5), Cm(3.3), Cm(13.0), Cm(12.5),
             fill=api_bg, line_color=api_color, radius=True)
    add_textbox(slide, "FastAPI 서버 구조",
                Cm(0.8), Cm(3.5), Cm(12.5), Cm(0.7),
                font_size=10, bold=True, color=api_color)

    server_layers = [
        ("api_app.py", "FastAPI 앱 생성 + /health 엔드포인트", api_color),
        ("routes.py", "POST /api/query/generate\nQueryGenerateRequest → response_model", api_color),
        ("query_service.py", "generate_stream_query(question, mode)\n→ StreamOrchestrator 위임", RGBColor(0x00, 0x52, 0x9B)),
        ("stream_orchestrator.py", "mode 선택 → natural/prompt/rag 실행", RGBColor(0x00, 0x46, 0x85)),
    ]
    for i, (name, desc, col) in enumerate(server_layers):
        flow_box(slide, name, Cm(0.8), Cm(4.3) + i * Cm(2.3),
                 w=Cm(12.5), h=Cm(1.0), fill=col, font_size=9)
        add_textbox(slide, desc,
                    Cm(0.9), Cm(5.4) + i * Cm(2.3), Cm(12.3), Cm(0.8),
                    font_size=8, color=C_GRAY)
        if i < len(server_layers) - 1:
            arrow_v(slide, Cm(7.05), Cm(4.3) + i * Cm(2.3) + Cm(1.0),
                    Cm(4.3) + (i + 1) * Cm(2.3), color=api_color)

    # mode alias 매핑
    add_rect(slide, Cm(0.8), Cm(13.0), Cm(12.5), Cm(2.5),
             fill=RGBColor(0xF8, 0xFD, 0xFF), line_color=api_color, radius=True)
    add_textbox(slide, "mode 별칭 매핑 (query_service.py)",
                Cm(1.0), Cm(13.2), Cm(12.0), Cm(0.6),
                font_size=8.5, bold=True, color=api_color)
    aliases = [
        "'natural'   → 'natural_llm'",
        "'prompt'    → 'prompt_llm'",
        "'rag_prompt'→ 'rag_prompt_llm'",
    ]
    for i, a in enumerate(aliases):
        add_textbox(slide, a, Cm(1.2), Cm(13.9) + i * Cm(0.55), Cm(12.0), Cm(0.5),
                    font_size=8, color=C_GRAY)

    # 우측: Request/Response 예시
    add_rect(slide, Cm(14.0), Cm(3.3), Cm(19.3), Cm(7.5),
             fill=RGBColor(0x1A, 0x1A, 0x2E), radius=True)
    add_textbox(slide, "Sample Request / Response",
                Cm(14.3), Cm(3.5), Cm(18.8), Cm(0.6),
                font_size=9, bold=True, color=C_ACCENT)
    req_text = (
        "POST /api/query/generate\n"
        "{\n"
        '  "question": "최근 30일 각 설비별 에러 건수는?",\n'
        '  "mode": "rag_prompt"\n'
        "}"
    )
    add_textbox(slide, req_text,
                Cm(14.3), Cm(4.2), Cm(18.8), Cm(3.0),
                font_size=8.5, color=RGBColor(0xA8, 0xFF, 0xC4))

    add_rect(slide, Cm(14.0), Cm(7.3), Cm(19.3), Cm(3.2),
             fill=RGBColor(0x12, 0x12, 0x20), radius=True)
    resp_text = (
        "HTTP 200 OK\n"
        "{\n"
        '  "success": true,\n'
        '  "mode": "rag_prompt",\n'
        '  "question": "최근 30일 각 설비별 에러 건수는?",\n'
        '  "generated_query": "SELECT EQPID, COUNT(*) ..."\n'
        "}"
    )
    add_textbox(slide, resp_text,
                Cm(14.3), Cm(7.5), Cm(18.8), Cm(2.8),
                font_size=8, color=RGBColor(0xFF, 0xE0, 0x82))

    # debug_client 설명
    add_rect(slide, Cm(14.0), Cm(11.0), Cm(19.3), Cm(4.8),
             fill=api_bg, line_color=api_color, radius=True)
    add_textbox(slide, "debug_client.py — 디버그 클라이언트",
                Cm(14.3), Cm(11.2), Cm(19.0), Cm(0.6),
                font_size=9, bold=True, color=api_color)
    debug_items = [
        "• 고정 질문을 서버에 POST — VS Code 브레이크포인트 확인용",
        "• STREAM_SERVER_ENDPOINT 환경변수로 대상 지정",
        "• STREAM_SERVER_MODE로 mode 설정 (기본: 'prompt')",
        "• python -m Root_Stream.server.debug_client 로 실행",
        "• response.json() 전체를 로거로 출력해 응답 확인",
        "• 현재 구현됨: GET /health + POST /api/query/generate",
        "• TODO: SQL 실행 결과 포함 엔드포인트 미구현",
    ]
    for i, item in enumerate(debug_items):
        add_textbox(slide, item, Cm(14.4), Cm(11.9) + i * Cm(0.56), Cm(18.8), Cm(0.52),
                    font_size=7.8, color=RGBColor(0x33, 0x33, 0x33))

    add_key_message(slide, "API_SERVER = FastAPI로 HTTP 인터페이스를 제공. mode 별칭으로 외부 클라이언트가 쉽게 접근하고, debug_client로 빠른 브레이크포인트 테스트가 가능하다.")
    add_notes(slide, "서버는 Root_Stream/server/api_app.py(FastAPI 앱) + routes.py(엔드포인트) + query_service.py(오케스트레이터 위임) 구조입니다. uvicorn으로 실행하고 debug_client.py로 테스트합니다.")


# ───────── 슬라이드 10: End-to-End 시퀀스 ─────────
def slide_10_e2e_sequence(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "End-to-End 시퀀스 다이어그램",
                  "사용자 질문 → Planner → Stream mode 분기 → SQL Guard → DB 실행 → 답변")

    # 액터 정의
    actors = [
        ("사용자", C_YELLOW),
        ("API/CLI", RGBColor(0x00, 0x96, 0x88)),
        ("Planner", RGBColor(0xFF, 0x8F, 0x00)),
        ("Orchestrator", C_BLUE1),
        ("EmbedSvc", C_INGEST),
        ("ChromaDB", RGBColor(0x00, 0x46, 0x85)),
        ("LLM", RGBColor(0x8B, 0x00, 0x00)),
        ("SQL Guard", RGBColor(0xE6, 0x37, 0x00)),
        ("MSSQL", RGBColor(0x55, 0x55, 0x55)),
    ]
    n = len(actors)
    margin_l = Cm(0.4)
    margin_r = Cm(0.4)
    usable_w = SLIDE_W - margin_l - margin_r
    col_gap = usable_w / n
    actor_y = Cm(3.0)
    actor_h = Cm(0.9)
    actor_w = Cm(3.4)
    life_y_start = actor_y + actor_h
    life_y_end = SLIDE_H - Cm(2.0)

    actor_centers = []
    for i, (name, col) in enumerate(actors):
        cx = margin_l + i * col_gap + col_gap / 2 - actor_w / 2
        flow_box(slide, name, cx, actor_y, w=actor_w, h=actor_h, fill=col, font_size=7.5)
        center_x = cx + actor_w / 2
        actor_centers.append(center_x)
        add_rect(slide, center_x - Cm(0.04), life_y_start, Cm(0.08),
                 life_y_end - life_y_start, fill=RGBColor(0xCC, 0xCC, 0xCC))

    # 메시지 시퀀스
    msgs = [
        (0, 1, Cm(4.2), "① 질문 입력 (question, mode, use_planner)", C_BLUE1),
        (1, 2, Cm(5.1), "② run_planner_step(question) [선택]", RGBColor(0xFF, 0x8F, 0x00)),
        (2, 1, Cm(6.0), "③ query_type 반환 (예: DB_THEN_RAG)", RGBColor(0xFF, 0x8F, 0x00)),
        (1, 3, Cm(6.9), "④ generate_stream_query(question, mode)", C_BLUE1),
        (3, 4, Cm(7.8), "⑤ embed_query(question) [RAG 모드]", C_INGEST),
        (4, 3, Cm(8.7), "⑥ query_embedding 반환", C_INGEST),
        (3, 5, Cm(9.6), "⑦ collection.query(embedding, top_k=3)", RGBColor(0x00, 0x46, 0x85)),
        (5, 3, Cm(10.5), "⑧ top-k chunks 반환", RGBColor(0x00, 0x46, 0x85)),
        (3, 6, Cm(11.4), "⑨ generate(system_prompt + context + question)", RGBColor(0x8B, 0x00, 0x00)),
        (6, 3, Cm(12.3), "⑩ 생성된 SQL 반환", RGBColor(0x8B, 0x00, 0x00)),
        (3, 7, Cm(13.2), "⑪ validate_query_sql(sql) — SELECT-only 검증", RGBColor(0xE6, 0x37, 0x00)),
        (7, 3, Cm(14.1), "⑫ 검증 통과 SQL 반환", RGBColor(0xE6, 0x37, 0x00)),
        (3, 8, Cm(15.0), "⑬ execute(sql) [--execute-sql 옵션 시]", RGBColor(0x55, 0x55, 0x55)),
        (8, 3, Cm(15.9), "⑭ rows 반환", RGBColor(0x55, 0x55, 0x55)),
        (3, 1, Cm(16.5), "⑮ StreamResult / E2E 최종 답변", C_BLUE1),
        (1, 0, Cm(17.1), "⑯ HTTP 응답 / CLI JSON 출력", C_YELLOW),
    ]

    for src, dst, yy, label, col in msgs:
        x1 = actor_centers[src]
        x2 = actor_centers[dst]
        add_arrow(slide, x1, yy, x2, yy, color=col, width=1)
        mx = (min(x1, x2) + max(x1, x2)) / 2
        add_textbox(slide, label,
                    mx - Cm(3.2), yy - Cm(0.38), Cm(6.4), Cm(0.38),
                    font_size=6, color=col, align=PP_ALIGN.CENTER)

    add_key_message(slide, "Planner(선택) → Stream mode 분기 → RAG 검색 → LLM → SQL Guard → DB 실행 → 답변. SQL Guard는 SELECT-only 안전 게이트이다.")
    add_notes(slide, "전체 시퀀스: Planner는 선택적 단계입니다. rag_prompt_llm 모드에서는 embed→ChromaDB 검색→LLM 경로를 거칩니다. SQL Guard가 항상 중간에 개입합니다. DB 실행은 --execute-sql 옵션 또는 E2E 함수 호출 시 실행됩니다.")


# ───────── 슬라이드 11: 실행 명령 / 운영 방법 ─────────
def slide_11_operations(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "실행 명령 / 운영 방법",
                  "CLI / SQL 실행 / RAG 사용 / FastAPI 서버 실행 예시")

    code_bg = RGBColor(0x1A, 0x1A, 0x2E)
    code_color = RGBColor(0xA8, 0xFF, 0xC4)
    comment_color = RGBColor(0x88, 0xBB, 0x88)
    label_color = C_BLUE1

    sections = [
        {
            "title": "1) CLI — natural_llm / prompt_llm",
            "code": (
                "# Root_Stream/config/config.yaml에서 mode 설정 후 실행\n"
                "python -m Root_Stream.main \\\n"
                "  --config Root_Stream/config/config.yaml \\\n"
                "  --question \"최근 에러 TOP 5를 보여줘\"\n"
                "\n"
                "# mode: natural_llm  →  스키마 없이 직접 LLM 호출\n"
                "# mode: prompt_llm   →  schema_context + business_rules 주입"
            ),
        },
        {
            "title": "2) CLI — SQL 실행 포함 (MSSQL 연결 필요)",
            "code": (
                "python -m Root_Stream.main \\\n"
                "  --question \"설비별 에러 건수를 조회해줘\" \\\n"
                "  --execute-sql\n"
                "\n"
                "# --execute-sql: 생성된 SQL을 MSSQL에서 실제 실행\n"
                "# config.yaml database.host / .database / .password 설정 필요"
            ),
        },
        {
            "title": "3) RAG 사용 (rag_prompt_llm)",
            "code": (
                "# config.yaml에서 mode: rag_prompt_llm 설정\n"
                "# retrieval.chroma_path: ../Root_Ingest/data/chroma\n"
                "# retrieval.collection_name: doc_chunks\n"
                "\n"
                "# Root_Ingest 실행 먼저 (ChromaDB 생성)\n"
                "python -m Root_Ingest.ingest.ingest_pipeline \\\n"
                "  --config Root_Ingest/config/config.yaml"
            ),
        },
        {
            "title": "4) FastAPI 서버 실행",
            "code": (
                "# 서버 시작\n"
                "uvicorn Root_Stream.server.api_app:app \\\n"
                "  --host 0.0.0.0 --port 8000 --reload\n"
                "\n"
                "# 헬스 체크\n"
                "curl http://localhost:8000/health\n"
                "\n"
                "# 디버그 클라이언트 실행\n"
                "python -m Root_Stream.server.debug_client \\\n"
                "  --mode rag_prompt --question \"에러 로그 조회\""
            ),
        },
    ]

    col_w = (SLIDE_W - Cm(1.4)) / 2 - Cm(0.3)
    positions = [
        (Cm(0.5), Cm(3.3)),
        (Cm(0.5 + col_w + 0.6), Cm(3.3)),
        (Cm(0.5), Cm(10.3)),
        (Cm(0.5 + col_w + 0.6), Cm(10.3)),
    ]

    for (gx, gy), sec in zip(positions, sections):
        box_h = Cm(6.7)
        add_rect(slide, gx, gy, col_w, Cm(0.75), fill=C_BLUE1)
        add_textbox(slide, sec["title"],
                    gx + Cm(0.2), gy + Cm(0.08), col_w - Cm(0.3), Cm(0.6),
                    font_size=9, bold=True, color=C_WHITE)
        add_rect(slide, gx, gy + Cm(0.75), col_w, Cm(5.8), fill=code_bg)
        add_textbox(slide, sec["code"],
                    gx + Cm(0.2), gy + Cm(0.9), col_w - Cm(0.3), Cm(5.5),
                    font_size=7.8, color=code_color)

    add_key_message(slide, "CLI, FastAPI, 노트북 모두 동일한 StreamOrchestrator를 재사용한다. 진입점만 다르다.")
    add_notes(slide, "Root_Stream.main이 CLI 진입점이고, Root_Stream.server.api_app이 FastAPI 진입점입니다. 둘 다 build_stream_orchestrator()를 호출합니다.")


# ───────── 슬라이드 12: 구현 범위 / 한계 / TODO ─────────
def slide_12_scope_todo(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "현재 구현 범위 / 한계 / TODO",
                  "소스 기준 실제 구현됨 vs 미구현 / 향후 확장 포인트")

    done_color = C_INGEST
    limit_color = C_GRAY
    todo_color = RGBColor(0x3A, 0x86, 0xFF)

    # ── 구현됨 ──
    add_rect(slide, Cm(0.5), Cm(3.2), Cm(10.3), Cm(13.5),
             fill=RGBColor(0xE8, 0xF7, 0xEE), line_color=done_color, radius=True)
    add_rect(slide, Cm(0.5), Cm(3.2), Cm(10.3), Cm(1.1), fill=done_color, radius=True)
    add_textbox(slide, "✅ 현재 구현됨",
                Cm(0.8), Cm(3.35), Cm(10.0), Cm(0.8),
                font_size=11, bold=True, color=C_WHITE)
    impl_items = [
        "Planner: LLM → JSON 파싱 → 검증 → PlannerPlan",
        "natural_llm: system_prompt + 질문 직접 LLM 전달",
        "prompt_llm: 템플릿 + schema_context + business_rules",
        "rag_prompt_llm: embed → ChromaDB → LLM",
        "SQL Guard: SELECT-only, 금지 키워드 14종 차단",
        "SqlExecutorService: MSSQL 실행 (pymssql)",
        "StreamOrchestrator: mode 분기 통합 관리",
        "FastAPI 서버: POST /api/query/generate + /health",
        "debug_client.py: 서버 브레이크포인트 테스트",
        "Root_Ingest 파이프라인: 수집→파싱→청킹→임베딩→ChromaDB",
        "E2E runner: run_planner/db/rag/final_answer_step()",
        "config.yaml 기반 완전한 설정 주입 구조",
        "단위 테스트: config_loader, SQL guard, mode_resolution",
        "노트북: 단계별 검증 (01~06, planner_flow)",
    ]
    for i, item in enumerate(impl_items):
        add_textbox(slide, f"  • {item}",
                    Cm(0.8), Cm(4.5) + i * Cm(0.82), Cm(10.0), Cm(0.75),
                    font_size=8, color=RGBColor(0x1B, 0x5E, 0x20))

    # ── 한계 / 미구현 ──
    add_rect(slide, Cm(11.3), Cm(3.2), Cm(10.3), Cm(6.4),
             fill=RGBColor(0xFF, 0xF3, 0xE0), line_color=limit_color, radius=True)
    add_rect(slide, Cm(11.3), Cm(3.2), Cm(10.3), Cm(1.1), fill=limit_color, radius=True)
    add_textbox(slide, "⚠️ 현재 미구현 / 제외됨",
                Cm(11.6), Cm(3.35), Cm(10.0), Cm(0.8),
                font_size=11, bold=True, color=C_WHITE)
    limit_items = [
        "Planner 독립 API 엔드포인트 없음",
        "streaming 응답 (SSE/WebSocket) 미구현",
        "SQL 실행 결과 포함 API 엔드포인트 없음",
        "Planner 결과 캐싱 없음",
        "임베딩 결과 캐싱 없음",
        "인증/인가 (API Key, JWT 등) 없음",
    ]
    for i, item in enumerate(limit_items):
        add_textbox(slide, f"  • {item}",
                    Cm(11.6), Cm(4.5) + i * Cm(0.82), Cm(10.0), Cm(0.75),
                    font_size=8, color=RGBColor(0x99, 0x60, 0x00))

    # ── TODO / 향후 확장 ──
    add_rect(slide, Cm(11.3), Cm(9.8), Cm(10.3), Cm(6.9),
             fill=RGBColor(0xE3, 0xF2, 0xFD), line_color=todo_color, radius=True)
    add_rect(slide, Cm(11.3), Cm(9.8), Cm(10.3), Cm(1.1), fill=todo_color, radius=True)
    add_textbox(slide, "🚀 향후 확장 포인트",
                Cm(11.6), Cm(9.95), Cm(10.0), Cm(0.8),
                font_size=11, bold=True, color=C_WHITE)
    todo_items = [
        "Planner API 엔드포인트 추가",
        "Streaming 응답 지원 (AsyncGenerator)",
        "SQL 실행 결과 포함 E2E API",
        "임베딩 / Planner 결과 Redis 캐싱",
        "Fine-tuned LLM으로 교체 지원",
        "피드백 루프: 잘못된 SQL 재학습 데이터 적재",
        "멀티 컬렉션 RAG 검색",
        "API Key 기반 인증 미들웨어",
    ]
    for i, item in enumerate(todo_items):
        add_textbox(slide, f"  • {item}",
                    Cm(11.6), Cm(11.1) + i * Cm(0.76), Cm(10.0), Cm(0.7),
                    font_size=8, color=todo_color)

    # ── 기술 스택 요약 ──
    add_rect(slide, Cm(22.1), Cm(3.2), Cm(11.2), Cm(13.5),
             fill=RGBColor(0xF8, 0xF8, 0xFF), line_color=C_BLUE1, radius=True)
    add_rect(slide, Cm(22.1), Cm(3.2), Cm(11.2), Cm(1.1), fill=C_BLUE1, radius=True)
    add_textbox(slide, "🔧 기술 스택",
                Cm(22.4), Cm(3.35), Cm(11.0), Cm(0.8),
                font_size=11, bold=True, color=C_WHITE)
    tech_stack = [
        ("언어", "Python 3.x"),
        ("LLM", "Ollama (qwen2.5:14b) / OpenAI GPT"),
        ("임베딩", "sentence-transformers multilingual"),
        ("VectorDB", "ChromaDB (PersistentClient)"),
        ("RDBMS", "MSSQL (pymssql)"),
        ("서버", "FastAPI + uvicorn"),
        ("설정", "YAML (config.yaml)"),
        ("로깅", "Python logging → data/logs/"),
        ("테스트", "pytest + 단위 테스트"),
        ("노트북", "Jupyter Notebook (01~06 + E2E)"),
        ("PPT", "python-pptx"),
    ]
    for i, (k, v) in enumerate(tech_stack):
        add_textbox(slide, f"  {k}:",
                    Cm(22.4), Cm(4.5) + i * Cm(0.82), Cm(3.0), Cm(0.75),
                    font_size=8, bold=True, color=C_BLUE1)
        add_textbox(slide, v,
                    Cm(25.5), Cm(4.5) + i * Cm(0.82), Cm(7.5), Cm(0.75),
                    font_size=8, color=C_GRAY)

    add_key_message(slide, "핵심 기능은 모두 구현됨. 스트리밍·인증·캐싱·독립 Planner API가 주요 확장 포인트이다.")
    add_notes(slide, "소스 기준으로 분석한 구현 범위입니다. 스트리밍, 인증, 캐싱은 아직 미구현 상태입니다.")


# ───────── 슬라이드 13: 결론 ─────────
def slide_13_conclusion(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "결론",
                  "핵심 포인트 요약 — Planner & Stream의 역할 차이 · RAG & API 위치")

    # 상단 2개 핵심 박스 (Planner vs Stream)
    left_w = Cm(15.5)
    right_w = Cm(15.5)

    add_rect(slide, Cm(0.5), Cm(3.3), left_w, Cm(6.5),
             fill=RGBColor(0xFF, 0xF3, 0xE0), line_color=RGBColor(0xFF, 0x8F, 0x00), radius=True)
    add_textbox(slide, "🧠 Planner의 역할",
                Cm(0.8), Cm(3.5), left_w - Cm(0.5), Cm(0.8),
                font_size=13, bold=True, color=RGBColor(0xFF, 0x8F, 0x00))
    planner_points = [
        "• 질문을 받아 query_type을 결정하는 라우터",
        "• DB_ONLY / RAG_ONLY / GENERAL 등 6가지 타입",
        "• 복합 질문(DB_THEN_RAG)을 steps 배열로 분해",
        "• LLM → JSON 계획 → 검증까지만 담당",
        "• DB 실행, RAG 검색, SQL 생성은 포함하지 않음",
        "• planner_flow_runner.py가 E2E 흐름을 연결",
    ]
    for i, s in enumerate(planner_points):
        add_textbox(slide, s, Cm(0.9), Cm(4.5) + i * Cm(0.8), left_w - Cm(0.8), Cm(0.72),
                    font_size=9, color=RGBColor(0x55, 0x40, 0x00))

    add_rect(slide, Cm(17.5), Cm(3.3), right_w, Cm(6.5),
             fill=C_LITE, line_color=C_BLUE1, radius=True)
    add_textbox(slide, "⚙️ Stream의 역할",
                Cm(17.8), Cm(3.5), right_w - Cm(0.5), Cm(0.8),
                font_size=13, bold=True, color=C_BLUE1)
    stream_points = [
        "• 실제 SQL 생성 & 응답 반환을 담당하는 실행 엔진",
        "• 3가지 mode: natural / prompt / rag_prompt_llm",
        "• SQL Guard로 DML/DDL 차단 (안전 게이트)",
        "• MSSQL 실행 포함 (선택적 --execute-sql)",
        "• FastAPI 서버와 CLI 두 가지 진입점 모두 지원",
        "• config.yaml 한 곳에서 모든 동작 제어",
    ]
    for i, s in enumerate(stream_points):
        add_textbox(slide, s, Cm(17.9), Cm(4.5) + i * Cm(0.8), right_w - Cm(0.8), Cm(0.72),
                    font_size=9, color=RGBColor(0x00, 0x30, 0x60))

    # 중간: RAG & API SERVER 위치
    add_rect(slide, Cm(0.5), Cm(10.1), SLIDE_W - Cm(1.0), Cm(3.5),
             fill=RGBColor(0xE8, 0xF4, 0xFF), line_color=C_ACCENT, radius=True)
    add_textbox(slide, "🔍 RAG & API SERVER의 위치",
                Cm(0.8), Cm(10.3), Cm(12.0), Cm(0.7),
                font_size=11, bold=True, color=C_BLUE1)
    loc_items = [
        ("RAG (rag_prompt_llm)", "Root_Stream 내부 mode — Root_Ingest가 만든 ChromaDB를 읽어 동적 컨텍스트 제공. 최고 정확도의 SQL 생성 경로."),
        ("API SERVER\n(api_app + routes)", "Root_Stream의 HTTP 인터페이스 — 동일한 StreamOrchestrator를 FastAPI로 감싼 것. CLI와 완전 동일한 로직을 HTTP로 노출."),
    ]
    for i, (k, v) in enumerate(loc_items):
        add_textbox(slide, k, Cm(0.8), Cm(11.1) + i * Cm(1.1), Cm(5.5), Cm(1.0),
                    font_size=9, bold=True, color=C_BLUE1)
        add_textbox(slide, v, Cm(6.5), Cm(11.1) + i * Cm(1.1), SLIDE_W - Cm(7.5), Cm(1.0),
                    font_size=9, color=C_GRAY)

    # 하단 3개 핵심 메시지 카드
    cards = [
        ("1️⃣ Planner = 라우터",
         "질문 유형을 분류해 DB/RAG/LLM\n어떤 경로를 쓸지 결정한다",
         RGBColor(0xFF, 0x8F, 0x00)),
        ("2️⃣ Stream = 실행 엔진",
         "mode에 따라 SQL 생성/RAG/LLM을\n실제로 실행하고 결과를 반환한다",
         C_BLUE1),
        ("3️⃣ RAG + ChromaDB = 정확도의 열쇠",
         "Root_Ingest 품질이 Root_Stream 답변\n품질을 직접 결정한다",
         C_INGEST),
    ]
    card_w = (SLIDE_W - Cm(2.2)) / 3 - Cm(0.3)
    for i, (title, msg, col) in enumerate(cards):
        gx = Cm(0.7) + i * (card_w + Cm(0.35))
        gy = Cm(14.0)
        add_rect(slide, gx, gy, card_w, Cm(3.8), fill=col, radius=True)
        add_textbox(slide, title,
                    gx + Cm(0.3), gy + Cm(0.3), card_w - Cm(0.5), Cm(0.8),
                    font_size=10, bold=True, color=C_WHITE)
        add_textbox(slide, msg,
                    gx + Cm(0.3), gy + Cm(1.2), card_w - Cm(0.5), Cm(2.2),
                    font_size=9, color=C_WHITE)

    add_key_message(slide, "Planner는 라우터, Stream은 실행 엔진, RAG는 정확도 향상 레이어, API Server는 HTTP 진입점이다.")
    add_notes(slide, "핵심 정리: Planner가 경로를 결정하고 Stream이 실행합니다. RAG는 rag_prompt_llm mode에서 동작하며 Root_Ingest 품질에 의존합니다. API Server는 StreamOrchestrator를 그대로 HTTP로 노출한 것입니다.")


# ──────────────────────────────────────────────
# 메인
# ──────────────────────────────────────────────

def build_pptx(output_path: str) -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_system_overview(prs)
    slide_03_folder_structure(prs)
    slide_04_common_structure(prs)
    slide_05_planner(prs)
    slide_06_natural_llm(prs)
    slide_07_prompt_llm(prs)
    slide_08_rag_prompt_llm(prs)
    slide_09_api_server(prs)
    slide_10_e2e_sequence(prs)
    slide_11_operations(prs)
    slide_12_scope_todo(prs)
    slide_13_conclusion(prs)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"저장 완료: {output_path}")
    print(f"총 {len(prs.slides)}장 슬라이드")


if __name__ == "__main__":
    OUT = r"C:\Users\김민한\Desktop\docs\LLM\DB_TO_LLM_Architecture.pptx"
    build_pptx(OUT)
